# open source
import pandas as pd
import pptx
import pptx.util

from tools.logger import logger
from tools.helpers.models import Path
from tools.helpers.data_manager.filepath_manager import (save_file,
                                                         handle_permission_error,
                                                         handle_bad_extension_error,
                                                         handle_file_not_found_error)


def _write_powerpoint(output_path, prs):
    try:
        prs.save(output_path)
    except PermissionError as err:
        handle_permission_error(err, _write_powerpoint, output_path, handle_read_only_error=True)
    return output_path


def export_img_to_powerpoint(img_paths, output_path=None):
    """Export images to a PowerPoint presentation.

    :param img_paths: list of images to save in PowerPoint slides
    :param output_path: path to save output file
    :return: path
    """
    # Create a new presentation
    prs = pptx.Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slides = [prs.slides.add_slide(blank_slide_layout) for _ in range(len(img_paths))]
    for slide, img_path in zip(slides, img_paths):
        slide.shapes.add_picture(img_path, pptx.util.Inches(0), pptx.util.Inches(0.90), width=pptx.util.Inches(10))
    # Save file
    output_path = save_file(output_path, extension='.pptx', overwrite='ask')
    _write_powerpoint(output_path, prs)
    logger.info("Successfully written '{}' PowerPoint file to disk.".format(output_path))
    return output_path


def _write_excel(path, dataframes, sheet_names, index=False, **kwargs):
    """

    :param path:
    :param dataframes: list of dataframes
    :param sheet_names: list of sheet names (same index as dataframes)
    :param index: if True, index names are exported
    :param kwargs: keyword arguments for pd.DataFrame.to_excel function
    :return:
    """
    if path is None or path.isnone:
        logger.warning('No path set to write data to Excel! No file has been created.')
        return None
    try:
        logger.debug("Trying to open file '{}'".format(path))
        with pd.ExcelWriter(path) as writer:
            logger.debug("File {} opened.".format(path))
            for df, sheet_name in zip(dataframes, sheet_names):
                if isinstance(df.columns, pd.MultiIndex):
                    index = True  # index must be True if MultiIndex columns, otherwise NotImplementedError is raised
                df.to_excel(writer, sheet_name=sheet_name, index=index, **kwargs)
                logger.debug("Sheet {} written.".format(sheet_name))
        logger.debug("File {} closed.".format(path))
    except PermissionError as err:
        logger.exception(err)
        path = handle_permission_error(err, func=_write_excel, path=path, args=[dataframes, sheet_names, index],
                                       kwargs=kwargs, change_path_func=save_file, handle_read_only_error=True)
    except ValueError as err:
        if str(err).startswith("No engine for filetype"):
            path = handle_bad_extension_error(err, func=_write_excel, path=path, args=[dataframes, sheet_names, index],
                                              kwargs=kwargs, change_path_func=save_file, extension=".xlsx")
        else:
            raise err
    except FileNotFoundError as err:
        logger.exception(err)
        path = handle_file_not_found_error(err, func=_write_excel, path=path, args=[dataframes, sheet_names, index],
                                           kwargs=kwargs, change_path_func=save_file)
    return path


def save_excel_file(path, dataframes, sheet_names=None, extension=".xlsx",
                    check_path=False, to_excel_kwargs=None, **options):
    """Write dataframes to disk.

    :param path: Excel file path  (to be checked before with 'save_file' function).
    :param dataframes: dataframe or list of dataframes
    :param sheet_names:  sheet name or ordered list of sheet names
    :param extension: Excel file extension
    :param check_path: check path with save_file
    :param to_excel_kwargs: keyword arguments for _write_excel function
    :param options: options for save_file function
    :return: path
    """
    if check_path:
        path = save_file(path, extension=extension, **options)
    if isinstance(dataframes, pd.DataFrame):
        dataframes = [dataframes]
    length = len(dataframes)
    if sheet_names is None:
        sheet_names = ["Sheet{}".format(i + 1) for i in range(length)]
    elif isinstance(sheet_names, str):
        sheet_names = [sheet_names]
    if len(sheet_names) != length:
        logger.warning("Invalid sheet names! Default sheet names will be used.")
    to_excel_kwargs = {} if to_excel_kwargs is None else to_excel_kwargs
    n_path = _write_excel(Path(path), dataframes, sheet_names, **to_excel_kwargs)
    return n_path
